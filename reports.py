"""
Enterprise Report Generation System
Professional export functionality for PDF, Excel, PowerPoint, and custom formats
"""

import os
import io
import json
from datetime import datetime, timedelta
from typing import Dict, List, Optional, Any, Tuple
from decimal import Decimal
import base64

# Core dependencies
from flask import Blueprint, request, jsonify, send_file, current_app
from flask_jwt_extended import jwt_required, get_jwt_identity
import pandas as pd
import numpy as np

# Report generation libraries
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor, Color
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, PageBreak
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.linecharts import HorizontalLineChart
from reportlab.graphics.charts.piecharts import Pie
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT

# Excel generation
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.chart import LineChart, BarChart, PieChart, Reference
from openpyxl.utils.dataframe import dataframe_to_rows

# PowerPoint generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# Image processing
import matplotlib.pyplot as plt
import seaborn as sns
from matplotlib.backends.backend_agg import FigureCanvasAgg
import plotly.graph_objects as go
import plotly.io as pio

from models import db, User, Calculation, SubscriptionTier
from auth import require_subscription

# Reports Blueprint
reports_bp = Blueprint('reports', __name__, url_prefix='/api/reports')

class ReportGenerator:
    """Enterprise report generation engine"""
    
    def __init__(self):
        self.brand_colors = {
            'primary': '#6366f1',
            'secondary': '#8b5cf6',
            'accent': '#06b6d4',
            'success': '#10b981',
            'warning': '#f59e0b',
            'error': '#ef4444',
            'dark': '#1f2937',
            'light': '#f8fafc'
        }
        
        # Set matplotlib style
        plt.style.use('seaborn-v0_8')
        sns.set_palette("husl")
    
    def generate_pdf_report(self, calculation: Calculation, user: User, 
                          include_charts: bool = True, 
                          include_details: bool = True) -> bytes:
        """Generate comprehensive PDF report"""
        
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            textColor=HexColor(self.brand_colors['primary']),
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            spaceAfter=20,
            textColor=HexColor(self.brand_colors['dark'])
        )
        
        # Header
        story.append(Paragraph("Infinex ROI Analysis Report", title_style))
        story.append(Spacer(1, 20))
        
        # Executive Summary
        story.append(Paragraph("Executive Summary", heading_style))
        
        summary_data = [
            ['Metric', 'Value'],
            ['Investment Amount', f"${calculation.investment_amount:,.2f}"],
            ['Projected ROI', f"{calculation.roi_percentage:.2f}%"],
            ['Payback Period', f"{calculation.payback_period_months} months"],
            ['Net Present Value', f"${calculation.npv:,.2f}" if calculation.npv else "N/A"],
            ['Risk Score', f"{calculation.risk_score:.1f}/10" if calculation.risk_score else "N/A"]
        ]
        
        summary_table = Table(summary_data, colWidths=[2*inch, 3*inch])
        summary_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), HexColor(self.brand_colors['primary'])),
            ('TEXTCOLOR', (0, 0), (-1, 0), HexColor('#ffffff')),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 14),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), HexColor('#f8fafc')),
            ('GRID', (0, 0), (-1, -1), 1, HexColor('#e5e7eb'))
        ]))
        
        story.append(summary_table)
        story.append(Spacer(1, 30))
        
        # Project Details
        if include_details:
            story.append(Paragraph("Project Details", heading_style))
            
            details_data = [
                ['Project Type', calculation.project_type.replace('_', ' ').title()],
                ['Company Size', calculation.company_size.title()],
                ['Target Industry', calculation.target_industry.replace('_', ' ').title()],
                ['Currency', calculation.currency],
                ['Timeline', f"{calculation.timeline_months} months"],
                ['Analysis Date', calculation.created_at.strftime('%B %d, %Y')]
            ]
            
            for label, value in details_data:
                story.append(Paragraph(f"<b>{label}:</b> {value}", styles['Normal']))
            
            story.append(Spacer(1, 30))
        
        # Financial Analysis
        story.append(Paragraph("Financial Analysis", heading_style))
        
        if calculation.results_data:
            results = calculation.results_data
            
            # ROI Analysis
            story.append(Paragraph("<b>Return on Investment Analysis</b>", styles['Heading3']))
            
            roi_text = f"""
            The projected ROI of {calculation.roi_percentage:.2f}% indicates a 
            {'strong' if calculation.roi_percentage > 20 else 'moderate' if calculation.roi_percentage > 10 else 'conservative'} 
            return potential. With an investment of ${calculation.investment_amount:,.2f}, 
            the projected revenue is ${calculation.projected_revenue:,.2f} over {calculation.timeline_months} months.
            """
            
            story.append(Paragraph(roi_text, styles['Normal']))
            story.append(Spacer(1, 20))
            
            # Risk Assessment
            if calculation.risk_score:
                story.append(Paragraph("<b>Risk Assessment</b>", styles['Heading3']))
                
                risk_level = "Low" if calculation.risk_score < 4 else "Medium" if calculation.risk_score < 7 else "High"
                risk_text = f"""
                The risk score of {calculation.risk_score:.1f}/10 indicates a {risk_level.lower()} risk profile. 
                This assessment considers market volatility, industry factors, and project complexity.
                """
                
                story.append(Paragraph(risk_text, styles['Normal']))
                story.append(Spacer(1, 20))
        
        # Charts
        if include_charts and calculation.results_data:
            story.append(PageBreak())
            story.append(Paragraph("Visual Analysis", heading_style))
            
            # Generate charts and add to PDF
            chart_images = self._generate_chart_images(calculation)
            
            for chart_name, chart_buffer in chart_images.items():
                chart_image = Image(chart_buffer, width=5*inch, height=3*inch)
                story.append(chart_image)
                story.append(Spacer(1, 20))
        
        # Footer
        story.append(PageBreak())
        footer_text = f"""
        <para align="center">
        <font size="10" color="{self.brand_colors['dark']}">
        Generated by Infinex ROI Calculator<br/>
        Report Date: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}<br/>
        User: {user.full_name} ({user.email})<br/>
        Calculation ID: {calculation.id}
        </font>
        </para>
        """
        
        story.append(Paragraph(footer_text, styles['Normal']))
        
        # Build PDF
        doc.build(story)
        
        buffer.seek(0)
        return buffer.getvalue()
    
    def generate_excel_report(self, calculation: Calculation, user: User,
                            include_data_tables: bool = True) -> bytes:
        """Generate comprehensive Excel report with multiple sheets"""
        
        buffer = io.BytesIO()
        wb = Workbook()
        
        # Remove default sheet
        wb.remove(wb.active)
        
        # Summary Sheet
        summary_ws = wb.create_sheet("Executive Summary")
        self._create_excel_summary_sheet(summary_ws, calculation, user)
        
        # Detailed Analysis Sheet
        if include_data_tables:
            analysis_ws = wb.create_sheet("Detailed Analysis")
            self._create_excel_analysis_sheet(analysis_ws, calculation)
        
        # Charts Sheet
        charts_ws = wb.create_sheet("Charts & Visualizations")
        self._create_excel_charts_sheet(charts_ws, calculation)
        
        # Raw Data Sheet
        if calculation.results_data:
            data_ws = wb.create_sheet("Raw Data")
            self._create_excel_data_sheet(data_ws, calculation)
        
        # Save workbook
        wb.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
    
    def generate_powerpoint_report(self, calculation: Calculation, user: User) -> bytes:
        """Generate professional PowerPoint presentation"""
        
        prs = Presentation()
        
        # Slide 1: Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "ROI Analysis Report"
        subtitle.text = f"Project: {calculation.project_type.replace('_', ' ').title()}\n" \
                       f"Prepared for: {user.company_name or user.full_name}\n" \
                       f"Date: {datetime.now().strftime('%B %d, %Y')}"
        
        # Slide 2: Executive Summary
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        slide.shapes.title.text = "Executive Summary"
        
        content = slide.placeholders[1].text_frame
        content.text = f"Investment Amount: ${calculation.investment_amount:,.2f}"
        
        p = content.add_paragraph()
        p.text = f"Projected ROI: {calculation.roi_percentage:.2f}%"
        
        p = content.add_paragraph()
        p.text = f"Payback Period: {calculation.payback_period_months} months"
        
        if calculation.npv:
            p = content.add_paragraph()
            p.text = f"Net Present Value: ${calculation.npv:,.2f}"
        
        # Slide 3: Financial Projections (with chart)
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        slide.shapes.title.text = "Financial Projections"
        
        # Add chart to slide
        if calculation.results_data:
            chart_data = CategoryChartData()
            chart_data.categories = ['Month 1', 'Month 6', 'Month 12', 'Month 18', 'Month 24']
            chart_data.add_series('ROI %', [5, 12, 18, 22, calculation.roi_percentage])
            
            x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(5)
            chart = slide.shapes.add_chart(
                prs.chart.XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
            ).chart
            
            chart.has_legend = True
            chart.legend.position = prs.chart.XL_LEGEND_POSITION.RIGHT
        
        # Slide 4: Risk Analysis
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = "Risk Analysis"
        
        if calculation.risk_score:
            risk_level = "Low" if calculation.risk_score < 4 else "Medium" if calculation.risk_score < 7 else "High"
            
            content = slide.placeholders[1].text_frame
            content.text = f"Risk Score: {calculation.risk_score:.1f}/10 ({risk_level} Risk)"
            
            p = content.add_paragraph()
            p.text = "Risk factors considered:"
            
            p = content.add_paragraph()
            p.text = "• Market volatility"
            p.level = 1
            
            p = content.add_paragraph()
            p.text = "• Industry trends"
            p.level = 1
            
            p = content.add_paragraph()
            p.text = "• Project complexity"
            p.level = 1
        
        # Save presentation
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
    
    def _generate_chart_images(self, calculation: Calculation) -> Dict[str, io.BytesIO]:
        """Generate chart images for PDF reports"""
        charts = {}
        
        if not calculation.results_data:
            return charts
        
        # ROI Progression Chart
        fig, ax = plt.subplots(figsize=(10, 6))
        
        months = list(range(1, calculation.timeline_months + 1))
        roi_progression = [
            (month / calculation.timeline_months) * calculation.roi_percentage 
            for month in months
        ]
        
        ax.plot(months, roi_progression, linewidth=3, color=self.brand_colors['primary'])
        ax.set_title('ROI Progression Over Time', fontsize=16, fontweight='bold')
        ax.set_xlabel('Months')
        ax.set_ylabel('ROI (%)')
        ax.grid(True, alpha=0.3)
        
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
        buffer.seek(0)
        charts['roi_progression'] = buffer
        plt.close()
        
        # Investment Breakdown Pie Chart
        if calculation.results_data.get('cost_analysis'):
            fig, ax = plt.subplots(figsize=(8, 8))
            
            cost_data = calculation.results_data['cost_analysis']
            labels = list(cost_data.keys())
            sizes = list(cost_data.values())
            
            colors = [self.brand_colors['primary'], self.brand_colors['secondary'], 
                     self.brand_colors['accent'], self.brand_colors['success']]
            
            ax.pie(sizes, labels=labels, autopct='%1.1f%%', colors=colors[:len(labels)])
            ax.set_title('Investment Breakdown', fontsize=16, fontweight='bold')
            
            buffer = io.BytesIO()
            plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
            buffer.seek(0)
            charts['investment_breakdown'] = buffer
            plt.close()
        
        return charts
    
    def _create_excel_summary_sheet(self, ws, calculation: Calculation, user: User):
        """Create Excel summary sheet"""
        
        # Title
        ws['A1'] = "Infinex ROI Analysis Report"
        ws['A1'].font = Font(size=20, bold=True, color="FFFFFF")
        ws['A1'].fill = PatternFill(start_color="6366F1", end_color="6366F1", fill_type="solid")
        ws.merge_cells('A1:D1')
        
        # User Information
        ws['A3'] = "Prepared for:"
        ws['B3'] = user.full_name
        ws['A4'] = "Company:"
        ws['B4'] = user.company_name or "N/A"
        ws['A5'] = "Date:"
        ws['B5'] = datetime.now().strftime('%B %d, %Y')
        
        # Key Metrics
        ws['A7'] = "Key Metrics"
        ws['A7'].font = Font(size=14, bold=True)
        
        metrics = [
            ("Investment Amount", f"${calculation.investment_amount:,.2f}"),
            ("Projected ROI", f"{calculation.roi_percentage:.2f}%"),
            ("Payback Period", f"{calculation.payback_period_months} months"),
            ("Net Present Value", f"${calculation.npv:,.2f}" if calculation.npv else "N/A"),
            ("Risk Score", f"{calculation.risk_score:.1f}/10" if calculation.risk_score else "N/A")
        ]
        
        for i, (metric, value) in enumerate(metrics, start=8):
            ws[f'A{i}'] = metric
            ws[f'B{i}'] = value
            ws[f'A{i}'].font = Font(bold=True)
        
        # Auto-adjust column widths
        for column in ws.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)
            ws.column_dimensions[column_letter].width = adjusted_width
    
    def _create_excel_analysis_sheet(self, ws, calculation: Calculation):
        """Create detailed analysis sheet"""
        
        ws['A1'] = "Detailed Financial Analysis"
        ws['A1'].font = Font(size=16, bold=True)
        
        if calculation.results_data:
            # Convert results to DataFrame for easy Excel export
            results_df = pd.DataFrame([calculation.results_data])
            
            for r in dataframe_to_rows(results_df, index=False, header=True):
                ws.append(r)
    
    def _create_excel_charts_sheet(self, ws, calculation: Calculation):
        """Create charts sheet in Excel"""
        
        ws['A1'] = "Charts & Visualizations"
        ws['A1'].font = Font(size=16, bold=True)
        
        # ROI Timeline Chart
        if calculation.timeline_months:
            months = list(range(1, calculation.timeline_months + 1))
            roi_values = [
                (month / calculation.timeline_months) * calculation.roi_percentage 
                for month in months
            ]
            
            # Add data
            ws['A3'] = "Month"
            ws['B3'] = "ROI %"
            
            for i, (month, roi) in enumerate(zip(months, roi_values), start=4):
                ws[f'A{i}'] = month
                ws[f'B{i}'] = roi
            
            # Create chart
            chart = LineChart()
            chart.title = "ROI Progression"
            chart.style = 2
            chart.x_axis.title = "Months"
            chart.y_axis.title = "ROI (%)"
            
            data = Reference(ws, min_col=2, min_row=3, max_row=3 + len(months))
            cats = Reference(ws, min_col=1, min_row=4, max_row=3 + len(months))
            chart.add_data(data, titles_from_data=True)
            chart.set_categories(cats)
            
            ws.add_chart(chart, "D3")
    
    def _create_excel_data_sheet(self, ws, calculation: Calculation):
        """Create raw data sheet"""
        
        ws['A1'] = "Raw Calculation Data"
        ws['A1'].font = Font(size=16, bold=True)
        
        # Input parameters
        ws['A3'] = "Input Parameters"
        ws['A3'].font = Font(size=12, bold=True)
        
        inputs = [
            ("Project Type", calculation.project_type),
            ("Company Size", calculation.company_size),
            ("Target Industry", calculation.target_industry),
            ("Investment Amount", calculation.investment_amount),
            ("Timeline (months)", calculation.timeline_months),
            ("Currency", calculation.currency)
        ]
        
        for i, (param, value) in enumerate(inputs, start=4):
            ws[f'A{i}'] = param
            ws[f'B{i}'] = str(value)
        
        # Results data
        if calculation.results_data:
            ws['A12'] = "Results Data"
            ws['A12'].font = Font(size=12, bold=True)
            
            row = 13
            for key, value in calculation.results_data.items():
                ws[f'A{row}'] = key
                ws[f'B{row}'] = str(value)
                row += 1

# Route handlers

@reports_bp.route('/calculation/<calculation_id>/pdf', methods=['GET'])
@require_subscription([SubscriptionTier.BASIC, SubscriptionTier.PROFESSIONAL, SubscriptionTier.ENTERPRISE])
def export_pdf_report(calculation_id):
    """Export calculation as PDF report"""
    try:
        user = request.current_user if hasattr(request, 'current_user') else None
        if not user:
            user_id = get_jwt_identity()
            user = User.query.get(user_id)
        
        calculation = Calculation.query.filter_by(
            id=calculation_id, 
            user_id=user.id
        ).first()
        
        if not calculation:
            return jsonify({
                'success': False,
                'message': 'Calculation not found'
            }), 404
        
        # Check user limits
        user_limits = user.get_limits()
        if 'pdf' not in user_limits.get('export_formats', []):
            return jsonify({
                'success': False,
                'message': 'PDF export not available in your subscription tier'
            }), 403
        
        # Generate report
        generator = ReportGenerator()
        
        include_charts = request.args.get('include_charts', 'true').lower() == 'true'
        include_details = request.args.get('include_details', 'true').lower() == 'true'
        
        pdf_data = generator.generate_pdf_report(
            calculation, user, include_charts, include_details
        )
        
        # Create response
        filename = f"roi_analysis_{calculation_id}_{datetime.now().strftime('%Y%m%d')}.pdf"
        
        return send_file(
            io.BytesIO(pdf_data),
            mimetype='application/pdf',
            as_attachment=True,
            download_name=filename
        )
        
    except Exception as e:
        current_app.logger.error(f"PDF export error: {e}")
        return jsonify({
            'success': False,
            'message': 'Failed to generate PDF report'
        }), 500

@reports_bp.route('/calculation/<calculation_id>/excel', methods=['GET'])
@require_subscription([SubscriptionTier.BASIC, SubscriptionTier.PROFESSIONAL, SubscriptionTier.ENTERPRISE])
def export_excel_report(calculation_id):
    """Export calculation as Excel report"""
    try:
        user = request.current_user if hasattr(request, 'current_user') else None
        if not user:
            user_id = get_jwt_identity()
            user = User.query.get(user_id)
        
        calculation = Calculation.query.filter_by(
            id=calculation_id, 
            user_id=user.id
        ).first()
        
        if not calculation:
            return jsonify({
                'success': False,
                'message': 'Calculation not found'
            }), 404
        
        # Check user limits
        user_limits = user.get_limits()
        if 'xlsx' not in user_limits.get('export_formats', []):
            return jsonify({
                'success': False,
                'message': 'Excel export not available in your subscription tier'
            }), 403
        
        # Generate report
        generator = ReportGenerator()
        
        include_data_tables = request.args.get('include_data', 'true').lower() == 'true'
        
        excel_data = generator.generate_excel_report(
            calculation, user, include_data_tables
        )
        
        # Create response
        filename = f"roi_analysis_{calculation_id}_{datetime.now().strftime('%Y%m%d')}.xlsx"
        
        return send_file(
            io.BytesIO(excel_data),
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True,
            download_name=filename
        )
        
    except Exception as e:
        current_app.logger.error(f"Excel export error: {e}")
        return jsonify({
            'success': False,
            'message': 'Failed to generate Excel report'
        }), 500

@reports_bp.route('/calculation/<calculation_id>/powerpoint', methods=['GET'])
@require_subscription([SubscriptionTier.PROFESSIONAL, SubscriptionTier.ENTERPRISE])
def export_powerpoint_report(calculation_id):
    """Export calculation as PowerPoint presentation"""
    try:
        user = request.current_user if hasattr(request, 'current_user') else None
        if not user:
            user_id = get_jwt_identity()
            user = User.query.get(user_id)
        
        calculation = Calculation.query.filter_by(
            id=calculation_id, 
            user_id=user.id
        ).first()
        
        if not calculation:
            return jsonify({
                'success': False,
                'message': 'Calculation not found'
            }), 404
        
        # Check user limits
        user_limits = user.get_limits()
        if 'pptx' not in user_limits.get('export_formats', []):
            return jsonify({
                'success': False,
                'message': 'PowerPoint export not available in your subscription tier'
            }), 403
        
        # Generate report
        generator = ReportGenerator()
        pptx_data = generator.generate_powerpoint_report(calculation, user)
        
        # Create response
        filename = f"roi_presentation_{calculation_id}_{datetime.now().strftime('%Y%m%d')}.pptx"
        
        return send_file(
            io.BytesIO(pptx_data),
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=filename
        )
        
    except Exception as e:
        current_app.logger.error(f"PowerPoint export error: {e}")
        return jsonify({
            'success': False,
            'message': 'Failed to generate PowerPoint presentation'
        }), 500

@reports_bp.route('/calculation/<calculation_id>/formats', methods=['GET'])
@jwt_required()
def get_available_formats(calculation_id):
    """Get available export formats for user's subscription"""
    try:
        user_id = get_jwt_identity()
        user = User.query.get(user_id)
        
        if not user:
            return jsonify({
                'success': False,
                'message': 'User not found'
            }), 404
        
        calculation = Calculation.query.filter_by(
            id=calculation_id, 
            user_id=user.id
        ).first()
        
        if not calculation:
            return jsonify({
                'success': False,
                'message': 'Calculation not found'
            }), 404
        
        user_limits = user.get_limits()
        available_formats = user_limits.get('export_formats', ['json'])
        
        format_info = {
            'json': {
                'name': 'JSON Data',
                'description': 'Raw calculation data in JSON format',
                'endpoint': f'/api/reports/calculation/{calculation_id}/json'
            },
            'csv': {
                'name': 'CSV Export',
                'description': 'Spreadsheet-compatible data export',
                'endpoint': f'/api/reports/calculation/{calculation_id}/csv'
            },
            'pdf': {
                'name': 'PDF Report',
                'description': 'Professional formatted report with charts',
                'endpoint': f'/api/reports/calculation/{calculation_id}/pdf'
            },
            'xlsx': {
                'name': 'Excel Workbook',
                'description': 'Multi-sheet Excel report with data and charts',
                'endpoint': f'/api/reports/calculation/{calculation_id}/excel'
            },
            'pptx': {
                'name': 'PowerPoint Presentation',
                'description': 'Executive presentation with key insights',
                'endpoint': f'/api/reports/calculation/{calculation_id}/powerpoint'
            }
        }
        
        return jsonify({
            'success': True,
            'data': {
                'available_formats': [
                    {**format_info[fmt], 'format': fmt} 
                    for fmt in available_formats if fmt in format_info
                ],
                'subscription_tier': user.subscription_tier.value,
                'calculation_id': calculation_id
            }
        }), 200
        
    except Exception as e:
        current_app.logger.error(f"Get formats error: {e}")
        return jsonify({
            'success': False,
            'message': 'Failed to get available formats'
        }), 500

# Export the blueprint
__all__ = ['reports_bp', 'ReportGenerator']